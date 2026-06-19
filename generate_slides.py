from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title_text, content_text):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = title_text
    tf = content.text_frame
    tf.text = content_text
    
    # Điều chỉnh font size cho dễ nhìn
    for p in tf.paragraphs:
        p.font.size = Pt(20)

# Slide 1
add_slide("Slide 1 – Giới thiệu đề tài", "• Tiêu đề: BA Agent sử dụng Agentic RAG và Wikipedia Knowledge Base\n• Bối cảnh dự án\n• Bài toán đặt ra\n• Mục tiêu giải quyết")

# Slide 2
add_slide("Slide 2 – Vấn đề hiện tại (Nỗi đau của BA)", "• BA phải đọc tài liệu thủ công -> Mất thời gian\n• Dễ bỏ sót yêu cầu\n• Khó chuyển đổi kiến thức thô thành User Story chuẩn")

# Slide 3
add_slide("Slide 3 – Mục tiêu hệ thống", "Hệ thống có thể làm gì?\n• Đọc hiểu PDF chuyên ngành\n• Hỏi làm rõ yêu cầu (Elicitation / Clarification)\n• Tra cứu Wikipedia động\n• Sinh User Story & Business Rules\n• Xuất CSV tương thích Jira")

# Slide 4
add_slide("Slide 4 – Kiến trúc tổng thể", "Sơ đồ luồng xử lý:\nUser\n ↓\nBA Agent\n ↓\nTool Selection (PDF RAG Tool / Wiki RAG Tool)\n ↓\nReasoning\n ↓\nStructured Output\n ↓\nCSV Export")

# Slide 5
add_slide("Slide 5 – Agent Workflow", "• Core: LangGraph ReAct Agent\n• Quy trình:\nQuestion\n ↓\nReasoning\n ↓\nTool Calling\n ↓\nObservation\n ↓\nFinal Answer")

# Slide 6
add_slide("Slide 6 – PDF RAG Pipeline", "Quy trình chuẩn bị dữ liệu nội bộ:\n\nPDF\n ↓\nChunking (Section-aware)\n ↓\nEmbedding\n ↓\nChromaDB\n ↓\nHybrid Search\n ↓\nRe-ranking")

# Slide 7
add_slide("Slide 7 – Wikipedia RAG Tool", "Không dùng Summary trực tiếp.\nQuy trình:\nWikipedia\n ↓\nContent\n ↓\nChunk\n ↓\nEmbedding\n ↓\nTemporary Vector Store\n ↓\nRetrieval")

# Slide 8
add_slide("Slide 8 – Clarification Mechanism", "Cơ chế chống Ảo giác (Hallucination):\n\nVí dụ Input: Xây dựng hệ thống khách sạn\n\nAgent phản ứng: \n- Có thanh toán online không?\n- Có những vai trò (Actor) nào?")

# Slide 9
add_slide("Slide 9 – Demo Results", "[Chèn Hình Ảnh / Video Demo tại đây]\n\n• Tool Calling\n• Clarification\n• User Story JSON/Markdown\n• CSV Export")

# Slide 10
add_slide("Slide 10 – Đánh giá", "Bảng kết quả Test:\n• PDF QA: Pass\n• Wiki Retrieval: Pass\n• Clarification: Pass\n• User Story: Pass\n• CSV Export: Pass")

# Slide 11
add_slide("Slide 11 – Hạn chế (Limitations)", "• Chưa tích hợp Jira API\n• Chưa có Traceability 1-1 (User Story map trực tiếp về trang PDF gốc)\n• Chưa có RAGAS Evaluation liên tục trên pipeline tự động\n• Chưa có Memory dài hạn (Long-term DB)")

# Slide 12
add_slide("Slide 12 – Hướng phát triển (Future Work)", "• Jira Integration\n• Confluence Integration\n• Traceability Mapping\n• Evaluation Framework\n• Multi-Agent Collaboration (Thêm Dev Agent, Tester Agent)")

prs.save("BA_Agent_Presentation.pptx")
print("Đã tạo file PPTX thành công!")
